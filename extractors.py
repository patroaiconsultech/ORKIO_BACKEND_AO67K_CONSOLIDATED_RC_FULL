from __future__ import annotations

from typing import Tuple
import io
import re
import zipfile
import xml.etree.ElementTree as ET

from pypdf import PdfReader
from docx import Document

MAX_EXTRACT_CHARS = 1_200_000  # safety guard


def _trim(text: str) -> str:
    text = (text or "").replace("\x00", "")
    if len(text) > MAX_EXTRACT_CHARS:
        text = text[:MAX_EXTRACT_CHARS]
    return text


def _xml_text(node: ET.Element) -> str:
    return "".join(node.itertext()).strip()


def _xlsx_shared_strings(zf: zipfile.ZipFile) -> list[str]:
    try:
        raw = zf.read("xl/sharedStrings.xml")
    except Exception:
        return []
    try:
        root = ET.fromstring(raw)
    except Exception:
        return []
    out: list[str] = []
    for item in root.findall(".//{*}si"):
        out.append(_xml_text(item))
    return out


def _xlsx_sheet_targets(zf: zipfile.ZipFile) -> list[tuple[str, str]]:
    try:
        workbook = ET.fromstring(zf.read("xl/workbook.xml"))
        rels = ET.fromstring(zf.read("xl/_rels/workbook.xml.rels"))
    except Exception:
        return []

    rel_targets: dict[str, str] = {}
    for rel in rels.findall(".//{*}Relationship"):
        rid = rel.attrib.get("Id") or ""
        target = rel.attrib.get("Target") or ""
        if not rid or not target:
            continue
        if not target.startswith("xl/"):
            target = "xl/" + target.lstrip("/")
        rel_targets[rid] = target

    sheets: list[tuple[str, str]] = []
    for sheet in workbook.findall(".//{*}sheet"):
        name = sheet.attrib.get("name") or "Sheet"
        rid = sheet.attrib.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id") or ""
        target = rel_targets.get(rid)
        if target:
            sheets.append((name, target))
    return sheets


def _xlsx_cell_value(cell: ET.Element, shared: list[str]) -> str:
    cell_type = cell.attrib.get("t") or ""
    if cell_type == "inlineStr":
        return _xml_text(cell)

    value_node = cell.find("{*}v")
    raw = (value_node.text if value_node is not None else "") or ""
    raw = raw.strip()
    if not raw:
        return ""

    if cell_type == "s":
        try:
            idx = int(raw)
            return shared[idx].strip() if 0 <= idx < len(shared) else ""
        except Exception:
            return ""
    if cell_type == "b":
        return "TRUE" if raw == "1" else "FALSE"
    return raw


def _xlsx_col_index(cell_ref: str, fallback: int) -> int:
    match = re.match(r"([A-Za-z]+)", cell_ref or "")
    if not match:
        return fallback
    col = 0
    for ch in match.group(1).upper():
        col = col * 26 + (ord(ch) - ord("A") + 1)
    return max(0, col - 1)


def _extract_xlsx_text(content: bytes) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore

        workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        parts: list[str] = []
        for sheet in workbook.worksheets:
            rows_out: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                if values:
                    rows_out.append(" | ".join(values))
            if rows_out:
                parts.append(f"Sheet: {sheet.title}\n" + "\n".join(rows_out))
        try:
            workbook.close()
        except Exception:
            pass
        return _trim("\n\n".join(parts))
    except Exception:
        pass

    parts: list[str] = []
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as zf:
            shared = _xlsx_shared_strings(zf)
            for sheet_name, target in _xlsx_sheet_targets(zf):
                try:
                    root = ET.fromstring(zf.read(target))
                except Exception:
                    continue
                rows_out: list[str] = []
                for row in root.findall(".//{*}sheetData/{*}row"):
                    values_by_col: dict[int, str] = {}
                    fallback_col = 0
                    for cell in row.findall("{*}c"):
                        value = _xlsx_cell_value(cell, shared)
                        if not value:
                            fallback_col += 1
                            continue
                        col_idx = _xlsx_col_index(cell.attrib.get("r") or "", fallback_col)
                        values_by_col[col_idx] = value
                        fallback_col = max(fallback_col + 1, col_idx + 1)
                    if values_by_col:
                        max_col = max(values_by_col)
                        row_values = [values_by_col.get(i, "") for i in range(max_col + 1)]
                        rows_out.append(" | ".join(v for v in row_values if v).strip())
                if rows_out:
                    parts.append(f"Sheet: {sheet_name}\n" + "\n".join(rows_out))
    except Exception:
        return ""
    return _trim("\n\n".join(parts))


def _extract_pptx_text(content: bytes) -> str:
    def _openxml_fallback() -> str:
        try:
            with zipfile.ZipFile(io.BytesIO(content)) as zf:
                slide_names = sorted(
                    [
                        name
                        for name in zf.namelist()
                        if re.match(r"ppt/slides/slide\d+\.xml$", name)
                    ],
                    key=lambda value: int(re.search(r"(\d+)", value).group(1)),
                )
                slides: list[str] = []
                for idx, slide_name in enumerate(slide_names, start=1):
                    try:
                        root = ET.fromstring(zf.read(slide_name))
                    except Exception:
                        continue
                    parts = [
                        "".join(node.itertext()).strip()
                        for node in root.findall(".//{*}t")
                        if "".join(node.itertext()).strip()
                    ]
                    if parts:
                        slides.append(f"Slide {idx}\n" + "\n".join(parts))
                return _trim("\n\n".join(slides))
        except Exception:
            return ""

    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        return _openxml_fallback()

    try:
        deck = Presentation(io.BytesIO(content))
    except Exception:
        return _openxml_fallback()

    slides: list[str] = []
    for idx, slide in enumerate(deck.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text:
                parts.append(str(text).strip())
            try:
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        values = [str(cell.text or "").strip() for cell in row.cells if str(cell.text or "").strip()]
                        if values:
                            parts.append(" | ".join(values))
            except Exception:
                continue
        if parts:
            slides.append(f"Slide {idx}\n" + "\n".join(parts))
    text = _trim("\n\n".join(slides))
    return text or _openxml_fallback()


def extract_text(filename: str, content: bytes) -> Tuple[str, int]:
    """
    Extract plain text from common office formats.
    Returns: (text, extracted_chars)
    """
    name = (filename or "").lower()

    # PDF
    if name.endswith(".pdf"):
        try:
            reader = PdfReader(io.BytesIO(content))
            parts = []
            for page in reader.pages:
                try:
                    t = page.extract_text() or ""
                except Exception:
                    t = ""
                if t:
                    parts.append(t)
            text = "\n\n".join(parts)
        except Exception:
            text = ""
        text = _trim(text)
        return text, len(text)

    # DOCX
    if name.endswith(".docx"):
        try:
            doc = Document(io.BytesIO(content))
            parts = [p.text for p in doc.paragraphs if p.text]
            text = "\n".join(parts)
        except Exception:
            text = ""
        text = _trim(text)
        return text, len(text)

    # XLSX
    if name.endswith(".xlsx"):
        text = _extract_xlsx_text(content)
        return text, len(text)

    # PPTX
    if name.endswith(".pptx"):
        text = _extract_pptx_text(content)
        return text, len(text)

    # TXT/MD/CSV fallback
    try:
        text = content.decode("utf-8", errors="ignore")
    except Exception:
        text = ""
    text = _trim(text)
    return text, len(text)

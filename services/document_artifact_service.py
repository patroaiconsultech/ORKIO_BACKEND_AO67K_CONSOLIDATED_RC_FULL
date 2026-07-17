from __future__ import annotations

import csv
import io
import multiprocessing as mp
import os
import re
import textwrap
import threading
from dataclasses import asdict, dataclass
from typing import Any, Dict, Iterable, List, Mapping, Optional


_ABSOLUTE_MAX_REQUEST_BYTES = 8 * 1024 * 1024
_ABSOLUTE_MAX_ROWS = 2_000
_ABSOLUTE_MAX_COLUMNS = 80
_ABSOLUTE_MAX_CELL_CHARS = 4_000
_ABSOLUTE_MAX_TOTAL_CELLS = 50_000
_ABSOLUTE_MAX_ARTIFACT_BYTES = 10 * 1024 * 1024
_ABSOLUTE_MAX_TIMEOUT_SECONDS = 60.0


class DocumentArtifactError(Exception):
    """Base DOCIO artifact error."""


class DocumentArtifactLimitError(DocumentArtifactError, ValueError):
    """Trusted generation limit was exceeded."""


class DocumentArtifactTimeoutError(DocumentArtifactError, TimeoutError):
    """Isolated generator exceeded its deadline and was terminated."""


class DocumentArtifactWorkerError(DocumentArtifactError, RuntimeError):
    """Isolated generator failed or exited without a valid result."""


class DocumentArtifactConcurrencyError(DocumentArtifactError, RuntimeError):
    """This replica has no free isolated generation slot."""


_GENERATION_STATE_LOCK = threading.Lock()
_GENERATION_ACTIVE = 0


@dataclass(frozen=True)
class DocumentArtifactLimits:
    max_request_bytes: int = 2 * 1024 * 1024
    max_rows: int = 1_000
    max_columns: int = 50
    max_cell_chars: int = 2_000
    max_total_cells: int = 25_000
    max_artifact_bytes: int = 5 * 1024 * 1024
    generation_timeout_seconds: float = 30.0

    @classmethod
    def from_env(cls) -> "DocumentArtifactLimits":
        return cls(
            max_request_bytes=_bounded_int(
                "ORKIO_DOCIO_MAX_REQUEST_BYTES",
                2 * 1024 * 1024,
                minimum=16 * 1024,
                maximum=_ABSOLUTE_MAX_REQUEST_BYTES,
            ),
            max_rows=_bounded_int(
                "ORKIO_DOCIO_MAX_ROWS",
                1_000,
                minimum=1,
                maximum=_ABSOLUTE_MAX_ROWS,
            ),
            max_columns=_bounded_int(
                "ORKIO_DOCIO_MAX_COLUMNS",
                50,
                minimum=1,
                maximum=_ABSOLUTE_MAX_COLUMNS,
            ),
            max_cell_chars=_bounded_int(
                "ORKIO_DOCIO_MAX_CELL_CHARS",
                2_000,
                minimum=1,
                maximum=_ABSOLUTE_MAX_CELL_CHARS,
            ),
            max_total_cells=_bounded_int(
                "ORKIO_DOCIO_MAX_TOTAL_CELLS",
                25_000,
                minimum=1,
                maximum=_ABSOLUTE_MAX_TOTAL_CELLS,
            ),
            max_artifact_bytes=_bounded_int(
                "ORKIO_DOCIO_MAX_ARTIFACT_BYTES",
                5 * 1024 * 1024,
                minimum=1_024,
                maximum=_ABSOLUTE_MAX_ARTIFACT_BYTES,
            ),
            generation_timeout_seconds=_bounded_float(
                "ORKIO_DOCIO_GENERATION_TIMEOUT_SECONDS",
                30.0,
                minimum=0.25,
                maximum=_ABSOLUTE_MAX_TIMEOUT_SECONDS,
            ),
        )


@dataclass
class GeneratedDocumentArtifact:
    filename: str
    content: bytes
    mime_type: str
    text_content: str
    format: str


def _bounded_int(name: str, default: int, *, minimum: int, maximum: int) -> int:
    raw = os.getenv(name)
    try:
        value = int(raw) if raw not in (None, "") else int(default)
    except (TypeError, ValueError):
        value = int(default)
    return max(minimum, min(value, maximum))


def _bounded_float(name: str, default: float, *, minimum: float, maximum: float) -> float:
    raw = os.getenv(name)
    try:
        value = float(raw) if raw not in (None, "") else float(default)
    except (TypeError, ValueError):
        value = float(default)
    return max(minimum, min(value, maximum))


def _try_acquire_generation_slot() -> bool:
    global _GENERATION_ACTIVE
    maximum = _bounded_int(
        "ORKIO_DOCIO_MAX_CONCURRENT_GENERATIONS",
        2,
        minimum=1,
        maximum=8,
    )
    with _GENERATION_STATE_LOCK:
        if _GENERATION_ACTIVE >= maximum:
            return False
        _GENERATION_ACTIVE += 1
        return True


def _release_generation_slot() -> None:
    global _GENERATION_ACTIVE
    with _GENERATION_STATE_LOCK:
        _GENERATION_ACTIVE = max(0, _GENERATION_ACTIVE - 1)


def _active_generation_slots() -> int:
    with _GENERATION_STATE_LOCK:
        return int(_GENERATION_ACTIVE)


def _safe_name(value: Any, fallback: str = "orkio_document") -> str:
    raw = _text(value).strip() or fallback
    raw = raw.replace("\r", " ").replace("\n", " ")
    raw = re.sub(r"[^A-Za-z0-9À-ÿ_. -]+", "_", raw)
    raw = re.sub(r"\s+", "_", raw).strip("._- ")
    return raw[:80] or fallback


def _text(value: Any) -> str:
    if value is None:
        return ""
    return str(value).replace("\x00", "").strip()


def neutralize_spreadsheet_formula(value: Any) -> str:
    """Keep untrusted spreadsheet content as text.

    Excel-compatible programs may execute cells beginning with =, +, - or @.
    Prefixing with an apostrophe is the portable safe default.
    """

    text = _text(value)
    if text.startswith(("=", "+", "-", "@")):
        return "'" + text
    return text


def _normalize_rows(
    rows: Optional[Iterable[Any]],
    limits: DocumentArtifactLimits,
) -> List[List[str]]:
    out: List[List[str]] = []
    total_cells = 0

    if rows is None:
        return out
    if isinstance(rows, (str, bytes, bytearray, Mapping)):
        raise DocumentArtifactLimitError("rows_must_be_an_array")

    for row_index, row in enumerate(rows):
        if row_index >= limits.max_rows:
            raise DocumentArtifactLimitError("max_rows_exceeded")

        if isinstance(row, Mapping):
            raw_values = list(row.values())
        elif isinstance(row, (list, tuple)):
            raw_values = list(row)
        else:
            raw_values = [row]

        if len(raw_values) > limits.max_columns:
            raise DocumentArtifactLimitError("max_columns_exceeded")

        total_cells += len(raw_values)
        if total_cells > limits.max_total_cells:
            raise DocumentArtifactLimitError("max_total_cells_exceeded")

        values: List[str] = []
        for value in raw_values:
            text = _text(value)
            if len(text) > limits.max_cell_chars:
                raise DocumentArtifactLimitError("max_cell_chars_exceeded")
            values.append(text)

        if any(value != "" for value in values):
            out.append(values)

    return out


def _normalize_slides(
    slides: Optional[Iterable[Any]],
    limits: DocumentArtifactLimits,
) -> List[Dict[str, Any]]:
    out: List[Dict[str, Any]] = []
    if slides is None:
        return out
    if isinstance(slides, (str, bytes, bytearray, Mapping)):
        raise DocumentArtifactLimitError("slides_must_be_an_array")

    for slide_index, slide in enumerate(slides):
        if slide_index >= 20:
            raise DocumentArtifactLimitError("max_slides_exceeded")
        if not isinstance(slide, Mapping):
            continue
        title = _text(slide.get("title"))[:120]
        raw_bullets = slide.get("bullets") or []
        if isinstance(raw_bullets, (str, bytes, bytearray, Mapping)):
            raw_bullets = [raw_bullets]
        bullets: List[str] = []
        for bullet in list(raw_bullets)[:10]:
            text = _text(bullet)[:320]
            if text and text not in bullets:
                bullets.append(text)
        if title or bullets:
            out.append(
                {
                    "source_slide": slide.get("source_slide") or slide_index + 1,
                    "title": title or f"Slide {slide_index + 1}",
                    "bullets": bullets,
                }
            )
    return out


def _normalize_source_plan(plan: Any) -> Optional[Dict[str, Any]]:
    if not isinstance(plan, Mapping):
        return None
    contract = _text(plan.get("contract")) or "PPTX_SOURCE_PLAN_V1"
    slides_payload = plan.get("slides") or []
    slides: List[Dict[str, Any]] = []
    if isinstance(slides_payload, list):
        for index, item in enumerate(slides_payload[:12], start=1):
            if not isinstance(item, Mapping):
                continue
            points = item.get("required_points") or []
            refs = item.get("source_refs") or []
            if isinstance(points, (str, bytes, bytearray, Mapping)):
                points = [points]
            if isinstance(refs, (str, bytes, bytearray, Mapping)):
                refs = [refs]
            slides.append(
                {
                    "index": int(item.get("index") or index),
                    "purpose": _shorten(item.get("purpose") or f"Slide {index}", 140),
                    "required_points": [_shorten(point, 180) for point in list(points)[:6]],
                    "source_refs": [_shorten(ref, 80) for ref in list(refs)[:6]],
                }
            )
    return {
        "contract": contract,
        "source_file_ids": [_shorten(value, 80) for value in list(plan.get("source_file_ids") or [])[:8]],
        "source_hashes": [_shorten(value, 80) for value in list(plan.get("source_hashes") or [])[:8]],
        "source_slide_count": int(plan.get("source_slide_count") or 0),
        "extracted_chars": int(plan.get("extracted_chars") or 0),
        "extracted_chunks": int(plan.get("extracted_chunks") or 0),
        "minimum_slide_count": int(plan.get("minimum_slide_count") or 0),
        "planned_slide_count": int(plan.get("planned_slide_count") or 0),
        "coverage_ratio": float(plan.get("coverage_ratio") or 0),
        "slides": slides,
        "unsupported_claims": [_shorten(value, 160) for value in list(plan.get("unsupported_claims") or [])[:8]],
        "premium_label_allowed": bool(plan.get("premium_label_allowed")),
    }


def _plain_table(rows: List[List[str]]) -> str:
    return "\n".join(
        " | ".join(cell for cell in row)
        for row in rows
        if any(cell != "" for cell in row)
    ).strip()


def _row_count(rows: List[List[str]]) -> int:
    if not rows:
        return 0
    return max(0, len(rows) - 1)


def _column_count(rows: List[List[str]]) -> int:
    if not rows:
        return 0
    return max(len(row) for row in rows)


def _body_paragraphs(body: str, *, max_items: int = 6) -> List[str]:
    paragraphs = [
        re.sub(r"\s+", " ", part).strip()
        for part in str(body or "").splitlines()
        if part.strip()
    ]
    if not paragraphs and body:
        paragraphs = [re.sub(r"\s+", " ", str(body or "")).strip()]
    return [paragraph[:700] for paragraph in paragraphs if paragraph][:max_items]


def _executive_summary(title: str, body: str, rows: List[List[str]]) -> List[str]:
    summary: List[str] = []
    count = _row_count(rows)
    cols = _column_count(rows)
    if count:
        summary.append(f"Base consolidada com {count} registro(s) e {cols} coluna(s).")
    if rows and rows[0]:
        summary.append("Campos principais: " + ", ".join(rows[0][:5]) + ".")
    for paragraph in _body_paragraphs(body, max_items=2):
        if paragraph not in summary:
            summary.append(paragraph)
    if not summary:
        summary.append(f"Artefato executivo gerado pela ORKIO para {title}.")
    return summary[:5]


def _source_quality_note(rows: List[List[str]]) -> str:
    if rows:
        return (
            "Fonte: contexto autorizado da conversa. "
            "Revise dados sensiveis antes de compartilhar externamente."
        )
    return "Fonte: conteudo gerado a partir da solicitacao explicita do usuario."


def _next_steps(rows: List[List[str]]) -> List[str]:
    if rows:
        return [
            "Validar a amostra contra a planilha original.",
            "Completar campos ausentes antes de decisoes comerciais.",
            "Priorizar proximas acoes por segmento, potencial e relacionamento.",
        ]
    return [
        "Revisar conteudo e ajustar tom ao publico-alvo.",
        "Adicionar dados reais quando disponiveis.",
        "Validar o artefato antes de uso externo.",
    ]


def _shorten(value: Any, limit: int = 160) -> str:
    text = re.sub(r"\s+", " ", _text(value)).strip()
    if len(text) <= limit:
        return text
    return text[: max(1, limit - 1)].rstrip() + "..."


def _premium_takeaway(title: str, body: str, rows: List[List[str]], slides: Optional[List[Dict[str, Any]]] = None) -> str:
    if slides:
        return f"{title}: narrativa preservada a partir de {len(slides)} slide(s) fonte."
    if rows:
        return f"{title}: amostra executiva com {_row_count(rows)} registro(s) e {_column_count(rows)} campo(s)."
    for paragraph in _body_paragraphs(body, max_items=1):
        return _shorten(paragraph, 120)
    return f"{title}: artefato executivo gerado com governanca."


def _source_quality_note_for(rows: List[List[str]], slides: Optional[List[Dict[str, Any]]] = None) -> str:
    if slides:
        return (
            "Fonte: outline extraido de PPTX autorizado vinculado a esta conversa. "
            "Revise identidade visual e imagens antes de uso externo."
        )
    return _source_quality_note(rows)


def _markdown(title: str, body: str, rows: List[List[str]]) -> str:
    parts = [
        f"# {title}".strip(),
        "",
        "## Resumo executivo",
        "",
        "\n".join(f"- {item}" for item in _executive_summary(title, body, rows)),
        "",
        "## Fonte e governanca",
        "",
        _source_quality_note(rows),
    ]
    body_paragraphs = _body_paragraphs(body)
    if body_paragraphs:
        parts.extend(["", "## Contexto", "", "\n\n".join(body_paragraphs)])
    if rows:
        parts.extend(["", "## Dados", "", _plain_table(rows)])
    parts.extend(["", "## Proximos passos", "", "\n".join(f"- {item}" for item in _next_steps(rows))])
    return "\n".join(part for part in parts if part is not None).strip() + "\n"


def _slides_markdown(title: str, slides: List[Dict[str, Any]]) -> str:
    parts = [f"# {title}".strip(), "", "## Resumo executivo", ""]
    parts.append(f"- Apresentacao reorganizada a partir de {len(slides)} slide(s) do PPTX fonte.")
    parts.append("- Conteudo preservado a partir do contexto autorizado da conversa.")
    for index, slide in enumerate(slides, start=1):
        parts.extend(["", f"## {index}. {slide.get('title') or 'Slide'}", ""])
        bullets = list(slide.get("bullets") or [])
        if bullets:
            parts.extend(f"- {bullet}" for bullet in bullets)
    parts.extend([
        "",
        "## Proximos passos",
        "",
        "- Validar narrativa, identidade visual e imagens antes de uso externo.",
    ])
    return "\n".join(str(part) for part in parts).strip() + "\n"


def _generate_md(
    title: str,
    body: str,
    rows: List[List[str]],
    basename: str,
) -> GeneratedDocumentArtifact:
    text = _markdown(title, body, rows)
    return GeneratedDocumentArtifact(
        filename=f"{basename}.md",
        content=text.encode("utf-8"),
        mime_type="text/markdown; charset=utf-8",
        text_content=text,
        format="md",
    )


def _generate_csv(
    title: str,
    body: str,
    rows: List[List[str]],
    basename: str,
) -> GeneratedDocumentArtifact:
    out = io.StringIO(newline="")
    writer = csv.writer(out)
    if rows:
        writer.writerows(
            [[neutralize_spreadsheet_formula(value) for value in row] for row in rows]
        )
    else:
        writer.writerow(["title", "content"])
        writer.writerow(
            [
                neutralize_spreadsheet_formula(title),
                neutralize_spreadsheet_formula(body),
            ]
        )
    text = out.getvalue()
    return GeneratedDocumentArtifact(
        filename=f"{basename}.csv",
        content=text.encode("utf-8-sig"),
        mime_type="text/csv; charset=utf-8",
        text_content=text,
        format="csv",
    )


def _generate_xlsx(
    title: str,
    body: str,
    rows: List[List[str]],
    basename: str,
) -> GeneratedDocumentArtifact:
    try:
        from openpyxl import Workbook  # type: ignore
        from openpyxl.styles import Alignment, Font, PatternFill  # type: ignore
        from openpyxl.utils import get_column_letter  # type: ignore
    except Exception as exc:
        raise RuntimeError("openpyxl_unavailable") from exc

    wb = Workbook()
    ws = wb.active
    ws.title = "Dados"
    ws.freeze_panes = "A2"
    if rows:
        for row in rows:
            ws.append([neutralize_spreadsheet_formula(value) for value in row])
        header_fill = PatternFill("solid", fgColor="1F4E78")
        header_font = Font(color="FFFFFF", bold=True)
        for cell in ws[1]:
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center", vertical="center")
        for column_cells in ws.columns:
            column = get_column_letter(column_cells[0].column)
            values = [str(cell.value or "") for cell in column_cells[:60]]
            ws.column_dimensions[column].width = max(12, min(42, max(len(v) for v in values) + 2))
        ws.auto_filter.ref = ws.dimensions
    else:
        ws.append(["Titulo", "Resumo"])
        ws.append([neutralize_spreadsheet_formula(title), neutralize_spreadsheet_formula(body)])

    meta = wb.create_sheet("Metadados")
    meta.append(["Campo", "Valor"])
    meta.append(["Titulo", neutralize_spreadsheet_formula(title)])
    meta.append(["Registros", str(_row_count(rows))])
    meta.append(["Colunas", str(_column_count(rows))])
    meta.append(["Fonte", _source_quality_note(rows)])
    meta.append(["Governanca", "Gerado sob solicitacao explicita do usuario"])
    for cell in meta[1]:
        cell.fill = PatternFill("solid", fgColor="404040")
        cell.font = Font(color="FFFFFF", bold=True)
    meta.column_dimensions["A"].width = 18
    meta.column_dimensions["B"].width = 72

    stream = io.BytesIO()
    wb.save(stream)
    text = _markdown(title, body, rows)
    return GeneratedDocumentArtifact(
        filename=f"{basename}.xlsx",
        content=stream.getvalue(),
        mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        text_content=text,
        format="xlsx",
    )


def _generate_docx(
    title: str,
    body: str,
    rows: List[List[str]],
    basename: str,
) -> GeneratedDocumentArtifact:
    try:
        from docx import Document  # type: ignore
        from docx.enum.text import WD_ALIGN_PARAGRAPH  # type: ignore
        from docx.shared import Inches, Pt, RGBColor  # type: ignore
    except Exception as exc:
        raise RuntimeError("python_docx_unavailable") from exc

    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.65)
    section.bottom_margin = Inches(0.65)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)

    styles = doc.styles
    styles["Normal"].font.name = "Arial"
    styles["Normal"].font.size = Pt(10)
    for style_name, size, color in (
        ("Title", 24, RGBColor(14, 17, 22)),
        ("Heading 1", 18, RGBColor(14, 17, 22)),
        ("Heading 2", 13, RGBColor(31, 78, 121)),
    ):
        style = styles[style_name]
        style.font.name = "Arial"
        style.font.size = Pt(size)
        style.font.color.rgb = color
        style.font.bold = True

    cover = doc.add_paragraph()
    cover.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = cover.add_run(title)
    run.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(14, 17, 22)

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle_run = subtitle.add_run(_premium_takeaway(title, body, rows))
    subtitle_run.font.size = Pt(11)
    subtitle_run.font.color.rgb = RGBColor(89, 89, 89)

    doc.add_heading("Resumo executivo", level=2)
    for item in _executive_summary(title, body, rows):
        doc.add_paragraph(item, style="List Bullet")
    doc.add_heading("Fonte e governanca", level=2)
    doc.add_paragraph(_source_quality_note(rows))
    body_paragraphs = _body_paragraphs(body)
    if body_paragraphs:
        doc.add_heading("Contexto", level=2)
    for paragraph in body_paragraphs:
        doc.add_paragraph(paragraph)
    if rows:
        doc.add_heading("Dados", level=2)
        width = max(len(row) for row in rows)
        table = doc.add_table(rows=0, cols=width)
        table.style = "Table Grid"
        for row in rows:
            cells = table.add_row().cells
            for index, value in enumerate(row):
                cells[index].text = value
        for cell in table.rows[0].cells:
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.bold = True
    doc.add_heading("Proximos passos", level=2)
    for item in _next_steps(rows):
        doc.add_paragraph(item, style="List Number")
    footer = section.footer.paragraphs[0]
    footer.text = "ORKIO | Artefato executivo gerado com governanca e fonte auditavel"
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER

    stream = io.BytesIO()
    doc.save(stream)
    text = _markdown(title, body, rows)
    return GeneratedDocumentArtifact(
        filename=f"{basename}.docx",
        content=stream.getvalue(),
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        text_content=text,
        format="docx",
    )


def _generate_pptx(
    title: str,
    body: str,
    rows: List[List[str]],
    basename: str,
    slides: Optional[List[Dict[str, Any]]] = None,
    source_plan: Optional[Dict[str, Any]] = None,
) -> GeneratedDocumentArtifact:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.dml.color import RGBColor  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore
        from pptx.enum.text import PP_ALIGN  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
    except Exception as exc:
        raise RuntimeError("python_pptx_unavailable") from exc

    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)

    bg = RGBColor(14, 17, 22)
    panel = RGBColor(28, 34, 45)
    accent = RGBColor(44, 123, 229)
    gold = RGBColor(235, 184, 78)
    white = RGBColor(245, 247, 250)
    muted = RGBColor(189, 197, 209)

    def _blank_slide() -> Any:
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = bg
        return slide

    def _textbox(slide: Any, text: str, left: float, top: float, width: float, height: float, *, size: int = 24, bold: bool = False, color: Any = None) -> Any:
        shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = shape.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = _shorten(text, 420)
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color or white
        return shape

    def _footer(slide: Any, index: int) -> None:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.08), Inches(13.333), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = accent
        line.line.fill.background()
        footer = _textbox(slide, f"ORKIO | Fonte auditavel | {index}", 0.55, 7.16, 8.0, 0.22, size=8, color=muted)
        footer.text_frame.paragraphs[0].font.bold = False

    def _bullet_panel(slide: Any, items: List[str], left: float, top: float, width: float, height: float, *, font_size: int = 19) -> None:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = panel
        box.line.color.rgb = RGBColor(52, 63, 82)
        tf = box.text_frame
        tf.clear()
        tf.margin_left = Inches(0.22)
        tf.margin_right = Inches(0.18)
        tf.margin_top = Inches(0.12)
        for index, item in enumerate(items[:7]):
            paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
            paragraph.text = _shorten(item, 145)
            paragraph.level = 0
            paragraph.font.name = "Arial"
            paragraph.font.size = Pt(font_size)
            paragraph.font.color.rgb = white if index == 0 else muted
            paragraph.font.bold = index == 0

    slide_no = 1
    title_slide = _blank_slide()
    side_bar = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.16), Inches(7.5))
    side_bar.fill.solid()
    side_bar.fill.fore_color.rgb = accent
    side_bar.line.fill.background()
    _textbox(title_slide, "ORKIO", 0.58, 0.45, 2.1, 0.32, size=15, bold=True, color=gold)
    _textbox(title_slide, _shorten(title, 92), 0.58, 1.35, 8.2, 1.25, size=40, bold=True, color=white)
    _textbox(title_slide, _premium_takeaway(title, body, rows, slides), 0.62, 2.75, 8.9, 0.6, size=18, color=muted)
    _bullet_panel(
        title_slide,
        [
            "Contrato premium: fonte real ou bloqueio transparente",
            "Entrega executiva com trilha de governanca",
            _source_quality_note_for(rows, slides),
        ],
        9.05,
        1.05,
        3.45,
        4.75,
        font_size=15,
    )
    _footer(title_slide, slide_no)

    slide_no += 1
    summary_slide = _blank_slide()
    _textbox(summary_slide, "Resumo executivo", 0.58, 0.48, 7.6, 0.5, size=30, bold=True, color=white)
    summary_items = _executive_summary(title, body, rows)
    if slides:
        summary_items = [
            f"{len(slides)} slide(s) fonte foram convertidos em narrativa executiva auditavel.",
            "A tese, os temas e os pontos centrais foram preservados como outline estruturado.",
            "A camada visual foi elevada, mas imagens/layout original ainda exigem validacao humana.",
        ]
    if source_plan:
        summary_items.append(
            "Contrato de fonte: "
            f"{source_plan.get('contract')} | "
            f"cobertura={source_plan.get('coverage_ratio')} | "
            f"slides planejados={source_plan.get('planned_slide_count')}."
        )
    _bullet_panel(summary_slide, summary_items, 0.7, 1.35, 11.8, 4.85, font_size=21)
    _footer(summary_slide, slide_no)

    if source_plan:
        slide_no += 1
        plan_slide = _blank_slide()
        _textbox(plan_slide, "Contrato de fonte e cobertura", 0.58, 0.48, 8.9, 0.5, size=30, bold=True, color=white)
        plan_items = [
            f"Contrato: {source_plan.get('contract')}",
            f"Fonte: {source_plan.get('source_slide_count')} bloco(s) | {source_plan.get('extracted_chars')} caractere(s) | {source_plan.get('extracted_chunks')} chunk(s)",
            f"Plano: minimo={source_plan.get('minimum_slide_count')} | planejado={source_plan.get('planned_slide_count')} | cobertura={source_plan.get('coverage_ratio')}",
            f"Premium label allowed: {str(bool(source_plan.get('premium_label_allowed'))).lower()}",
        ]
        file_ids = list(source_plan.get("source_file_ids") or [])
        if file_ids:
            plan_items.append("Source file IDs: " + ", ".join(file_ids[:3]))
        _bullet_panel(plan_slide, plan_items, 0.7, 1.25, 11.8, 2.25, font_size=17)
        refs = []
        for item in list(source_plan.get("slides") or [])[:5]:
            ref = ", ".join(item.get("source_refs") or []) or "source_ref:missing"
            refs.append(f"{item.get('index')}. {item.get('purpose')} | {ref}")
        _bullet_panel(plan_slide, refs or ["Sem source_refs suficientes para rotulo premium."], 0.7, 3.75, 11.8, 2.5, font_size=14)
        _footer(plan_slide, slide_no)

    if slides:
        for item in slides[:10]:
            slide_no += 1
            slide = _blank_slide()
            source_slide = item.get("source_slide") or slide_no - 2
            _textbox(slide, f"Slide fonte {source_slide}", 0.58, 0.45, 2.5, 0.26, size=12, bold=True, color=gold)
            _textbox(slide, str(item.get("title") or "Slide"), 0.58, 0.9, 9.8, 0.85, size=30, bold=True, color=white)
            bullets = list(item.get("bullets") or [])[:7] or ["Conteudo preservado a partir do slide fonte."]
            _bullet_panel(slide, bullets, 0.7, 2.0, 11.75, 4.35, font_size=18)
            _footer(slide, slide_no)
    elif rows:
        slide_no += 1
        slide = _blank_slide()
        _textbox(slide, "Dados selecionados da fonte", 0.58, 0.48, 8.8, 0.5, size=30, bold=True, color=white)
        table_rows = rows[: min(len(rows), 8)]
        table_cols = max(len(row) for row in table_rows)
        table = slide.shapes.add_table(
            len(table_rows),
            table_cols,
            Inches(0.55),
            Inches(1.35),
            Inches(12.25),
            Inches(4.95),
        ).table
        for row_index, row in enumerate(table_rows):
            for column_index, value in enumerate(row):
                cell = table.cell(row_index, column_index)
                cell.text = _shorten(value, 95)
                cell.fill.solid()
                cell.fill.fore_color.rgb = accent if row_index == 0 else panel
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = "Arial"
                    paragraph.font.size = Pt(10 if row_index else 11)
                    paragraph.alignment = PP_ALIGN.CENTER if row_index == 0 else PP_ALIGN.LEFT
                    paragraph.font.bold = row_index == 0
                    paragraph.font.color.rgb = white
        _footer(slide, slide_no)

    slide_no += 1
    next_slide = _blank_slide()
    _textbox(next_slide, "Proximos passos recomendados", 0.58, 0.48, 9.5, 0.5, size=30, bold=True, color=white)
    next_items = _next_steps(rows)
    if slides:
        next_items = [
            "Validar se a nova versao preserva a tese central do PPTX fonte.",
            "Aplicar identidade visual, imagens e layout final antes de apresentacao externa.",
            "Ajustar profundidade conforme publico-alvo, tempo disponivel e objetivo comercial.",
        ]
    _bullet_panel(next_slide, next_items, 0.7, 1.35, 11.8, 4.85, font_size=22)
    _footer(next_slide, slide_no)

    stream = io.BytesIO()
    deck.save(stream)
    text = _slides_markdown(title, slides) if slides else _markdown(title, body, rows)
    if source_plan:
        text += (
            "\n\n## PPTX_SOURCE_PLAN_V1\n\n"
            f"- contract: {source_plan.get('contract')}\n"
            f"- source_file_ids: {', '.join(source_plan.get('source_file_ids') or [])}\n"
            f"- source_slide_count: {source_plan.get('source_slide_count')}\n"
            f"- extracted_chars: {source_plan.get('extracted_chars')}\n"
            f"- extracted_chunks: {source_plan.get('extracted_chunks')}\n"
            f"- minimum_slide_count: {source_plan.get('minimum_slide_count')}\n"
            f"- planned_slide_count: {source_plan.get('planned_slide_count')}\n"
            f"- coverage_ratio: {source_plan.get('coverage_ratio')}\n"
            f"- premium_label_allowed: {str(bool(source_plan.get('premium_label_allowed'))).lower()}\n"
        )
    return GeneratedDocumentArtifact(
        filename=f"{basename}.pptx",
        content=stream.getvalue(),
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        text_content=text,
        format="pptx",
    )


def _generate_pdf(
    title: str,
    body: str,
    rows: List[List[str]],
    basename: str,
) -> GeneratedDocumentArtifact:
    try:
        from reportlab.lib.pagesizes import A4  # type: ignore
        from reportlab.pdfgen import canvas  # type: ignore
        from reportlab.pdfbase.pdfmetrics import stringWidth  # type: ignore
    except Exception as exc:
        raise RuntimeError("reportlab_unavailable") from exc

    stream = io.BytesIO()
    pdf = canvas.Canvas(stream, pagesize=A4)
    width, height = A4
    left = 72
    right = width - 72
    y = height - 96

    def draw_wrapped(value: str, *, font: str, size: int, leading: int) -> None:
        nonlocal y
        pdf.setFont(font, size)
        max_width = max(80, right - left)
        for source_line in (value or "").splitlines() or [""]:
            words = source_line.split()
            if not words:
                y -= leading
                continue
            current = ""
            for word in words:
                candidate = word if not current else f"{current} {word}"
                if stringWidth(candidate, font, size) <= max_width:
                    current = candidate
                    continue
                if y < 72:
                    pdf.showPage()
                    y = height - 72
                    pdf.setFont(font, size)
                pdf.drawString(left, y, current)
                y -= leading
                current = word
            if current:
                if y < 72:
                    pdf.showPage()
                    y = height - 72
                    pdf.setFont(font, size)
                pdf.drawString(left, y, current)
                y -= leading

    pdf.setFillColorRGB(0.055, 0.067, 0.086)
    pdf.rect(0, height - 74, width, 74, stroke=0, fill=1)
    pdf.setFillColorRGB(0.173, 0.482, 0.898)
    pdf.rect(0, height - 74, 12, 74, stroke=0, fill=1)
    pdf.setFillColorRGB(1, 1, 1)
    pdf.setFont("Helvetica-Bold", 18)
    pdf.drawString(left, height - 42, _shorten(title, 72))
    pdf.setFont("Helvetica", 8)
    pdf.setFillColorRGB(0.74, 0.77, 0.82)
    pdf.drawString(left, height - 58, "ORKIO | Artefato executivo com governanca e fonte auditavel")
    pdf.setFillColorRGB(0, 0, 0)

    draw_wrapped("Resumo executivo", font="Helvetica-Bold", size=13, leading=18)
    for item in _executive_summary(title, body, rows):
        draw_wrapped(f"- {item}", font="Helvetica", size=10, leading=15)
    y -= 4
    draw_wrapped("Fonte e governanca", font="Helvetica-Bold", size=13, leading=18)
    draw_wrapped(_source_quality_note(rows), font="Helvetica", size=10, leading=15)
    body_paragraphs = _body_paragraphs(body)
    if body_paragraphs:
        y -= 4
        draw_wrapped("Contexto", font="Helvetica-Bold", size=13, leading=18)
    for paragraph in body_paragraphs:
        draw_wrapped(paragraph, font="Helvetica", size=10, leading=15)
    if rows:
        y -= 6
        draw_wrapped("Dados", font="Helvetica-Bold", size=13, leading=18)
        draw_wrapped(_plain_table(rows), font="Helvetica", size=9, leading=13)
    y -= 6
    draw_wrapped("Proximos passos", font="Helvetica-Bold", size=13, leading=18)
    for item in _next_steps(rows):
        draw_wrapped(f"- {item}", font="Helvetica", size=10, leading=15)
    pdf.setFont("Helvetica", 8)
    pdf.setFillColorRGB(0.45, 0.45, 0.45)
    pdf.drawCentredString(width / 2, 36, "ORKIO | Validar dados sensiveis antes de uso externo")
    pdf.save()

    text = _markdown(title, body, rows)
    return GeneratedDocumentArtifact(
        filename=f"{basename}.pdf",
        content=stream.getvalue(),
        mime_type="application/pdf",
        text_content=text,
        format="pdf",
    )


def generate_document_artifact(
    payload: Dict[str, Any],
    *,
    limits: Optional[DocumentArtifactLimits] = None,
) -> GeneratedDocumentArtifact:
    trusted = limits or DocumentArtifactLimits.from_env()
    fmt = _text(payload.get("format") or payload.get("file_format") or "md").lower().lstrip(".")
    aliases = {
        "markdown": "md",
        "excel": "xlsx",
        "word": "docx",
        "powerpoint": "pptx",
    }
    fmt = aliases.get(fmt, fmt)

    title = _text(payload.get("title")) or "Documento Orkio"
    body = _text(payload.get("content") or payload.get("body") or payload.get("markdown"))
    rows = _normalize_rows(payload.get("rows") or payload.get("table_rows"), trusted)
    slides = _normalize_slides(payload.get("slides"), trusted)
    source_plan = _normalize_source_plan(payload.get("source_plan"))
    basename = _safe_name(payload.get("filename") or title)

    if len(title) > 180:
        raise DocumentArtifactLimitError("title_too_long")
    if len(body) > 200_000:
        raise DocumentArtifactLimitError("content_too_long")

    generators = {
        "md": _generate_md,
        "csv": _generate_csv,
        "xlsx": _generate_xlsx,
        "docx": _generate_docx,
        "pptx": _generate_pptx,
        "pdf": _generate_pdf,
    }
    generator = generators.get(fmt)
    if generator is None:
        raise ValueError(f"unsupported_document_format:{fmt}")

    if fmt == "pptx":
        artifact = _generate_pptx(title, body, rows, basename, slides=slides, source_plan=source_plan)
    else:
        artifact = generator(title, body, rows, basename)
    if len(artifact.content) > trusted.max_artifact_bytes:
        raise DocumentArtifactLimitError("max_artifact_bytes_exceeded")
    return artifact


def _artifact_worker(
    payload: Dict[str, Any],
    limits_payload: Dict[str, Any],
    sender: Any,
) -> None:
    try:
        artifact = generate_document_artifact(
            payload,
            limits=DocumentArtifactLimits(**limits_payload),
        )
        sender.send(
            {
                "ok": True,
                "artifact": {
                    "filename": artifact.filename,
                    "content": artifact.content,
                    "mime_type": artifact.mime_type,
                    "text_content": artifact.text_content,
                    "format": artifact.format,
                },
            }
        )
    except BaseException as exc:  # child boundary must report all failures
        try:
            sender.send(
                {
                    "ok": False,
                    "error_type": exc.__class__.__name__,
                    "error": str(exc)[:500],
                }
            )
        except Exception:
            pass
    finally:
        try:
            sender.close()
        except Exception:
            pass


def _terminate_process(process: mp.Process) -> None:
    if process.is_alive():
        process.terminate()
        process.join(timeout=2.0)
    if process.is_alive() and hasattr(process, "kill"):
        process.kill()
        process.join(timeout=2.0)
    if process.is_alive():
        raise DocumentArtifactWorkerError("generation_worker_could_not_be_terminated")


def generate_document_artifact_isolated(
    payload: Dict[str, Any],
    *,
    limits: Optional[DocumentArtifactLimits] = None,
    timeout_seconds: Optional[float] = None,
) -> GeneratedDocumentArtifact:
    """Generate in a terminable process guarded by a per-replica hard slot.

    A timeout always terminates and joins the worker before this function returns.
    The concurrency slot is released only after the worker is confirmed stopped,
    so a timed-out process cannot silently escape the resource boundary.
    """

    if not _try_acquire_generation_slot():
        raise DocumentArtifactConcurrencyError(
            "document_generation_concurrency_limit_reached"
        )

    trusted = limits or DocumentArtifactLimits.from_env()
    timeout = (
        trusted.generation_timeout_seconds
        if timeout_seconds is None
        else max(0.05, min(float(timeout_seconds), _ABSOLUTE_MAX_TIMEOUT_SECONDS))
    )

    preferred_method = os.getenv(
        "ORKIO_DOCIO_PROCESS_START_METHOD",
        "spawn",
    ).strip().lower()
    if preferred_method not in {"spawn", "fork", "forkserver"}:
        preferred_method = "spawn"
    if preferred_method not in mp.get_all_start_methods():
        preferred_method = "spawn"

    process: Optional[mp.Process] = None
    receiver = None
    sender = None
    try:
        context = mp.get_context(preferred_method)
        receiver, sender = context.Pipe(duplex=False)
        process = context.Process(
            target=_artifact_worker,
            args=(dict(payload), asdict(trusted), sender),
            daemon=True,
            name="orkio-docio-generator",
        )
        process.start()
        sender.close()
        sender = None

        if not receiver.poll(timeout):
            _terminate_process(process)
            raise DocumentArtifactTimeoutError("document_generation_timeout")

        result = receiver.recv()
        process.join(timeout=2.0)
        if process.is_alive():
            _terminate_process(process)
            raise DocumentArtifactWorkerError("generation_worker_did_not_exit")

        if not isinstance(result, dict) or not result.get("ok"):
            error_type = str((result or {}).get("error_type") or "WorkerError")
            error = str((result or {}).get("error") or "document_generation_failed")
            if error_type == "DocumentArtifactLimitError":
                raise DocumentArtifactLimitError(error)
            if error_type in {"ValueError"}:
                raise ValueError(error)
            if error_type in {"RuntimeError"}:
                raise DocumentArtifactWorkerError(error)
            raise DocumentArtifactWorkerError(f"{error_type}:{error}")

        artifact = result.get("artifact") or {}
        return GeneratedDocumentArtifact(
            filename=str(artifact["filename"]),
            content=bytes(artifact["content"]),
            mime_type=str(artifact["mime_type"]),
            text_content=str(artifact["text_content"]),
            format=str(artifact["format"]),
        )
    except EOFError as exc:
        if process is not None:
            _terminate_process(process)
        raise DocumentArtifactWorkerError(
            "generation_worker_exited_without_result"
        ) from exc
    finally:
        if receiver is not None:
            try:
                receiver.close()
            except Exception:
                pass
        if sender is not None:
            try:
                sender.close()
            except Exception:
                pass
        if process is not None and process.is_alive():
            _terminate_process(process)
        _release_generation_slot()
